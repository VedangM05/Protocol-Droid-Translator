import os
import tempfile
import openpyxl
from docx import Document
from pptx import Presentation
import fitz  # PyMuPDF (we use this to avoid needing poppler for pdf2image)
from PIL import Image
import shutil
import platform
import pytesseract

def get_installed_languages():
    """Returns a list of languages installed in Tesseract."""
    if not TESSERACT_PATH:
        return []
    try:
        return pytesseract.get_languages()
    except Exception:
        return []

# --- Tesseract Discovery (Windows) ---
def discover_tesseract():
    # 1. Check if it's already in PATH
    path_in_path = shutil.which("tesseract")
    if path_in_path:
        return path_in_path
        
    # 2. Check common Windows installation paths
    if platform.system() == "Windows":
        common_paths = [
            r'C:\Program Files\Tesseract-OCR\tesseract.exe',
            r'C:\Program Files (x86)\Tesseract-OCR\tesseract.exe',
            os.path.join(os.environ.get('LOCALAPPDATA', ''), 'Tesseract-OCR', 'tesseract.exe'),
            os.path.join(os.environ.get('USERPROFILE', ''), 'AppData', 'Local', 'Tesseract-OCR', 'tesseract.exe'),
        ]
        for p in common_paths:
            if os.path.exists(p):
                return p
    return None

TESSERACT_PATH = discover_tesseract()
if TESSERACT_PATH:
    pytesseract.pytesseract.tesseract_cmd = TESSERACT_PATH

def check_ffmpeg():
    """Checks if FFmpeg is installed and accessible."""
    return shutil.which("ffmpeg") is not None

class DocumentTranslator:
    def __init__(self, translation_engine):
        self.engine = translation_engine

    def _translate_text(self, text, src_lang, tgt_lang):
        if not text or not text.strip():
            return text
        try:
            return self.engine.translate(text=text, source_lang=src_lang, target_lang=tgt_lang)
        except Exception as e:
            print(f"Failed to translate segment: {text[:30]}... Error: {e}")
            return text

    def _apply_translations(self, operations, src_lang, tgt_lang):
        if not operations:
            return
            
        texts = [op["text"] for op in operations]
        
        # Batch translate via engine for significant speedup
        if hasattr(self.engine, 'translate_batch'):
            # Chunk into reasonable sizes to not overwhelm memory
            batch_size = 16 
            translated_texts = []
            for i in range(0, len(texts), batch_size):
                batch = texts[i:i+batch_size]
                try:
                    mode = getattr(self.engine, 'current_mode', 'Fast')
                    precision = getattr(self.engine, 'current_precision', 'Standard')
                    translated_batch = self.engine.translate_batch(batch, src_lang, tgt_lang, mode=mode, precision=precision)
                    translated_texts.extend(translated_batch)
                except Exception as e:
                    print(f"Batch translation failed: {e}. Falling back to single.")
                    for text in batch:
                        translated_texts.append(self._translate_text(text, src_lang, tgt_lang))
        else:
            translated_texts = [self._translate_text(t, src_lang, tgt_lang) for t in texts]
            
        for op, translated in zip(operations, translated_texts):
            # Extract text if returned as dictionary by translation engine
            text_result = translated["text"] if isinstance(translated, dict) else translated
            
            if text_result and text_result.strip():
                op["update"](text_result)

    def translate_docx(self, file_path, src_lang, tgt_lang):
        """Translates a Word document (.docx)."""
        doc = Document(file_path)
        operations = []
        
        # Translate paragraphs
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                def update_p(t, p=paragraph):
                    if p.runs:
                        p.runs[0].text = t
                        for i in range(1, len(p.runs)):
                            p.runs[i].text = ""
                operations.append({
                    "text": paragraph.text,
                    "update": update_p
                })

        # Translate tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        def update_c(t, c=cell):
                            c.text = t
                        operations.append({
                            "text": cell.text,
                            "update": update_c
                        })
                        
        self._apply_translations(operations, src_lang, tgt_lang)
        
        output_path = os.path.join(tempfile.gettempdir(), f"translated_{os.path.basename(file_path)}")
        doc.save(output_path)
        return output_path

    def translate_xlsx(self, file_path, src_lang, tgt_lang):
        """Translates an Excel document (.xlsx)."""
        workbook = openpyxl.load_workbook(file_path)
        operations = []
        
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value and isinstance(cell.value, str) and cell.value.strip():
                        # Exclude formulas
                        if not cell.value.startswith('='):
                            def update_c(t, c=cell):
                                c.value = t
                            operations.append({
                                "text": cell.value,
                                "update": update_c
                            })
                            
        self._apply_translations(operations, src_lang, tgt_lang)                            
        
        output_path = os.path.join(tempfile.gettempdir(), f"translated_{os.path.basename(file_path)}")
        workbook.save(output_path)
        return output_path

    def translate_pptx(self, file_path, src_lang, tgt_lang):
        """Translates a PowerPoint document (.pptx)."""
        prs = Presentation(file_path)
        operations = []
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        full_text = "".join([run.text for run in paragraph.runs])
                        if full_text.strip():
                            def update_p(t, p=paragraph):
                                if p.runs:
                                    p.runs[0].text = t
                                    for i in range(1, len(p.runs)):
                                        p.runs[i].text = ""
                            operations.append({
                                "text": full_text,
                                "update": update_p
                            })
                
                # Check for tables
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text_frame:
                                for paragraph in cell.text_frame.paragraphs:
                                    full_text = "".join([run.text for run in paragraph.runs])
                                    if full_text.strip():
                                        def update_p(t, p=paragraph):
                                            if p.runs:
                                                p.runs[0].text = t
                                                for i in range(1, len(p.runs)):
                                                    p.runs[i].text = ""
                                        operations.append({
                                            "text": full_text,
                                            "update": update_p
                                        })

        self._apply_translations(operations, src_lang, tgt_lang)
        
        output_path = os.path.join(tempfile.gettempdir(), f"translated_{os.path.basename(file_path)}")
        prs.save(output_path)
        return output_path

    def translate_pdf(self, file_path, src_lang, tgt_lang):
        """
        Translates a PDF document (.pdf).
        Uses OCR (Tesseract) because many government/local PDFs lack proper
        Unicode font dictionaries (causing mojibake/garbled text on standard extraction).
        """
        operations = []
        warning_msg = ""
        
        try:
            # Check for missing languages
            installed_langs = get_installed_languages()
            requested_langs = []
            missing_packs = []
            
            # Smart language selection logic
            for lang in ['mar', 'hin', 'eng']:
                if lang in installed_langs:
                    requested_langs.append(lang)
                elif lang != 'eng': # English is usually default
                    missing_packs.append(lang)
            
            if missing_packs and ('mar' in missing_packs or 'hin' in missing_packs):
                warning_msg = (
                    "⚠️  CRITICAL: MISSING LANGUAGE PACKS\n"
                    "Garbled text ('rubbish') detected? You're missing Indic support in Tesseract.\n\n"
                    "1. Download traineddata files:\n"
                    "   - Marathi: https://github.com/tesseract-ocr/tessdata_best/raw/main/mar.traineddata\n"
                    "   - Hindi: https://github.com/tesseract-ocr/tessdata_best/raw/main/hin.traineddata\n"
                    "2. Move them to: C:\\Program Files\\Tesseract-OCR\\tessdata\\\n"
                    "3. Restart the application.\n"
                    "--------------------------------------------------\n\n"
                )
            
            lang_str = "+".join(requested_langs) if requested_langs else "eng"
            
            doc = fitz.open(file_path)
            translated_pages = [""] * len(doc)
            
            # Use PyMuPDF to render pages to images
            for i in range(len(doc)):
                page = doc.load_page(i)
                # Render at 2.0x zoom (300DPI approx) - balance between speed and accuracy
                zoom = 2.0 
                mat = fitz.Matrix(zoom, zoom)
                pix = page.get_pixmap(matrix=mat)
                
                # Convert PyMuPDF pixmap to PIL Image
                if pix.alpha:
                    img = Image.frombytes("RGBA", [pix.width, pix.height], pix.samples)
                else:
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                
                # --- Advanced "Extreme OCR" Preprocessing ---
                from PIL import ImageEnhance, ImageFilter
                
                # 1. Sharpness injection to define characters
                img = img.filter(ImageFilter.SHARPEN)
                
                # 2. Convert to Grayscale
                img = img.convert('L') 
                
                # 3. Aggressive Contrast Boost
                enhancer = ImageEnhance.Contrast(img)
                img = enhancer.enhance(2.5)
                
                # 4. Binarization (Thresholding to pure Black & White)
                # This removes scanner noise and background textures common in legal docs
                # Point-thresholding is highly effective for Indic script definition
                img = img.point(lambda x: 255 if x > 128 else 0, mode='1')
                
                # Extract text using OCR with optimized configuration
                # Only use installed languages to prevent garbage character generation
                custom_config = r'--oem 3 --psm 3'
                text = pytesseract.image_to_string(img, lang=lang_str, config=custom_config)
                
                if text and text.strip():
                    def update_page_text(translated_text, idx=i):
                        translated_pages[idx] = translated_text
                        
                    operations.append({
                        "text": text.strip(),
                        "update": update_page_text
                    })
                    
            doc.close()

            # Process translation in batches for speed via our engine
            self._apply_translations(operations, src_lang, tgt_lang)
            
        except Exception as e:
            msg = (
                "⚠️ OCR ENGINE ERROR\n\n"
                "Please ensure Tesseract is installed correctly.\n\n"
                "1. Download: https://github.com/UB-Mannheim/tesseract/wiki\n"
                "2. Install with default settings.\n"
                "3. Restart BhashaEngine.\n\n"
                f"Technical Error: {e}"
            )
            translated_pages = [msg]
        
        # Since we can't easily recreate the PDF layout from scratch in Python reliably 
        # without heavy dependencies like reportlab padding, we return it as a structured .txt file
        output_path = os.path.join(tempfile.gettempdir(), f"translated_{os.path.basename(file_path)}.txt")
        with open(output_path, 'w', encoding='utf-8') as f:
            if warning_msg:
                f.write(warning_msg)
                
            for i, p_text in enumerate(translated_pages):
                if p_text:
                    f.write(f"--- Page {i+1} ---\n\n")
                    f.write(p_text)
                    f.write("\n\n")
                    
        return output_path
