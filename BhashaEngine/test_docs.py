import os
import tempfile
from docx import Document
import openpyxl
from pptx import Presentation
from translation.document_translator import DocumentTranslator

class MockTranslationEngine:
    def translate(self, text, source_lang, target_lang, domain="General"):
        # Just append something to mock translation without loading NLLB
        return f"[Translated]: {text}"

def create_dummy_docx(path):
    doc = Document()
    doc.add_paragraph("Hello world. This is a test document.")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Row 1 Col 1"
    table.cell(0, 1).text = "Row 1 Col 2"
    table.cell(1, 0).text = "Row 2 Col 1"
    table.cell(1, 1).text = "Row 2 Col 2"
    doc.save(path)

def create_dummy_xlsx(path):
    wb = openpyxl.Workbook()
    ws = wb.active
    ws['A1'] = "Excel Test Document"
    ws['B1'] = "Data 1"
    ws['A2'] = "Row 2"
    ws['B2'] = "Data 2"
    wb.save(path)

def create_dummy_pptx(path):
    prs = Presentation()
    slide_layout = prs.slide_layouts[1] # Title and content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Presentation Test"
    content.text = "This is a bullet point test."
    prs.save(path)

def main():
    test_dir = tempfile.gettempdir()
    docx_path = os.path.join(test_dir, "test.docx")
    xlsx_path = os.path.join(test_dir, "test.xlsx")
    pptx_path = os.path.join(test_dir, "test.pptx")
    
    print("Creating dummy documents...")
    create_dummy_docx(docx_path)
    create_dummy_xlsx(xlsx_path)
    create_dummy_pptx(pptx_path)
    
    engine = MockTranslationEngine()
    translator = DocumentTranslator(engine)
    
    print("Translating DOCX...")
    out_docx = translator.translate_docx(docx_path, "English", "Hindi")
    print(f"DOCX translated successfully to: {out_docx}")
    
    print("Translating XLSX...")
    out_xlsx = translator.translate_xlsx(xlsx_path, "English", "Hindi")
    print(f"XLSX translated successfully to: {out_xlsx}")
    
    print("Translating PPTX...")
    out_pptx = translator.translate_pptx(pptx_path, "English", "Hindi")
    print(f"PPTX translated successfully to: {out_pptx}")

if __name__ == "__main__":
    main()
